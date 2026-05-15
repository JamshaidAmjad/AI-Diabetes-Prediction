"""
Professional 15-slide PowerPoint — AI Diabetes Prediction Capstone
Author: Jamshaid Amjad | May 2026
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

def OxmlElement(tag):
    return etree.SubElement(etree.Element("_"), tag) if False else etree.Element(qn(tag) if ":" in tag else tag)

BASE = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
FIGS = os.path.join(BASE, "outputs", "figures")
OUT  = os.path.join(BASE, "outputs", "reports")
os.makedirs(OUT, exist_ok=True)

# ── Palette ───────────────────────────────────────────────────────────────────
NAVY        = RGBColor(0x1F, 0x38, 0x64)
BLUE        = RGBColor(0x2F, 0x75, 0xB6)
LIGHT_BLUE  = RGBColor(0xDE, 0xEA, 0xF1)
PALE_BLUE   = RGBColor(0xF0, 0xF5, 0xFA)
ORANGE      = RGBColor(0xED, 0x7D, 0x31)
GREEN       = RGBColor(0x70, 0xAD, 0x47)
DARK_GREEN  = RGBColor(0x37, 0x5C, 0x23)
RED         = RGBColor(0xC0, 0x00, 0x00)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
NEAR_BLACK  = RGBColor(0x1A, 0x1A, 0x1A)
DARK_GRAY   = RGBColor(0x40, 0x40, 0x40)
MED_GRAY    = RGBColor(0x70, 0x70, 0x70)
LIGHT_GRAY  = RGBColor(0xF5, 0xF5, 0xF5)
RULE_GRAY   = RGBColor(0xD0, 0xD0, 0xD0)

# ── Helpers ───────────────────────────────────────────────────────────────────
def _prs():
    p = Presentation()
    p.slide_width  = Inches(13.33)
    p.slide_height = Inches(7.5)
    return p

def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def R(slide, l, t, w, h, fill, line_color=None, line_w=0.75):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line_color:
        s.line.color.rgb = line_color
        s.line.width = Pt(line_w)
    else:
        s.line.fill.background()
    return s

def T(slide, text, l, t, w, h, size=13, bold=False, color=DARK_GRAY,
      align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = True
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return tb

def TM(slide, lines, l, t, w, h, size=12, bold=False, color=DARK_GRAY,
       align=PP_ALIGN.LEFT, spacing=5):
    """Multi-paragraph textbox."""
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = True
    tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(spacing)
        r = p.add_run(); r.text = line
        r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = color

def PIC(slide, fname, l, t, w=None, h=None):
    path = os.path.join(FIGS, fname)
    if not os.path.exists(path): return None
    kw = {}
    if w: kw['width']  = Inches(w)
    if h: kw['height'] = Inches(h)
    return slide.shapes.add_picture(path, Inches(l), Inches(t), **kw)

def cell_shade(cell, hex6):
    tc = cell._tc; p = tc.get_or_add_tcPr()
    shd = OxmlElement("w:shd")
    shd.set(qn("w:val"), "clear")
    shd.set(qn("w:color"), "auto")
    shd.set(qn("w:fill"), hex6)
    p.append(shd)

# ── Slide chrome ──────────────────────────────────────────────────────────────
def header(slide, title, subtitle=None):
    R(slide, 0, 0, 13.33, 1.12, NAVY)
    R(slide, 0, 0, 0.07,  1.12, ORANGE)          # left accent stripe
    T(slide, title, 0.22, 0.08, 12.8, 0.66,
      size=27, bold=True, color=WHITE)
    if subtitle:
        T(slide, subtitle, 0.22, 0.74, 12.8, 0.34,
          size=12, italic=True, color=LIGHT_BLUE)

def footer(slide, n):
    R(slide, 0, 7.12, 13.33, 0.38, NAVY)
    R(slide, 0, 7.12, 0.07,  0.38, ORANGE)
    T(slide, f"Jamshaid Amjad  |  AI in Healthcare — Capstone Project  |  May 2026   [{n}/15]",
      0.22, 7.15, 12.9, 0.32, size=9, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

# ── Section box helpers ───────────────────────────────────────────────────────
def card(slide, l, t, w, h, bg=LIGHT_BLUE, stripe=BLUE):
    R(slide, l, t,     w,    h,    bg)
    R(slide, l, t,     w,    0.055, stripe)

def card_title(slide, text, l, t, w, color=NAVY):
    T(slide, text, l+0.14, t+0.10, w-0.28, 0.42, size=14, bold=True, color=color)

def stat_card(slide, value, label, l, t, w=2.9, h=1.45,
              bg=LIGHT_BLUE, val_color=NAVY, stripe=ORANGE):
    R(slide, l, t, w, h, bg)
    R(slide, l, t, w, 0.055, stripe)
    T(slide, value, l+0.1, t+0.17, w-0.2, 0.72,
      size=32, bold=True, color=val_color, align=PP_ALIGN.CENTER)
    T(slide, label, l+0.1, t+0.88, w-0.2, 0.5,
      size=11, color=MED_GRAY, align=PP_ALIGN.CENTER)

def bullet_card(slide, title, items, l, t, w, h,
                bg=LIGHT_BLUE, stripe=BLUE, icon="▸", size=12):
    card(slide, l, t, w, h, bg, stripe)
    card_title(slide, title, l, t, w)
    row_h = min(0.46, (h - 0.56) / max(len(items), 1))
    for i, item in enumerate(items):
        T(slide, f"{icon}  {item}",
          l+0.14, t+0.55+i*row_h, w-0.28, row_h+0.05,
          size=size, color=DARK_GRAY)

def step_box(slide, num, title, desc, l, t, w=1.95, h=2.1):
    R(slide, l, t, w, h, LIGHT_BLUE)
    R(slide, l, t, w, 0.055, BLUE)
    R(slide, l, t+0.055, w, 0.5, BLUE)
    T(slide, num,   l+0.08, t+0.07, 0.35, 0.42, size=18, bold=True, color=WHITE)
    T(slide, title, l+0.42, t+0.07, w-0.55, 0.42, size=11, bold=True, color=WHITE)
    T(slide, desc,  l+0.1,  t+0.63, w-0.2, h-0.73, size=10, color=DARK_GRAY)

def metric_highlight(slide, label, value, l, t, w=2.2, h=0.95,
                     bg=LIGHT_BLUE, val_color=NAVY, stripe=BLUE):
    R(slide, l, t, w, h, bg)
    R(slide, l, t, 0.055, h, stripe)
    T(slide, label, l+0.16, t+0.08, w-0.25, 0.38, size=10, color=MED_GRAY)
    T(slide, value, l+0.16, t+0.44, w-0.25, 0.42, size=18, bold=True, color=val_color)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════════════════════
def s01_title(prs):
    sl = _blank(prs)
    # Background
    R(sl, 0, 0, 13.33, 7.5, NAVY)
    R(sl, 0, 0, 13.33, 4.5, RGBColor(0x17, 0x2B, 0x52))   # darker top half
    # Accent stripe
    R(sl, 0, 4.3, 13.33, 0.08, ORANGE)
    # Diagonal decorative block
    R(sl, 10.5, 0, 2.83, 4.35, RGBColor(0x25, 0x45, 0x78))

    T(sl, "AI IN HEALTHCARE", 0.6, 0.55, 9.0, 0.55,
      size=13, bold=True, color=ORANGE, italic=False)
    T(sl, "Early Diabetes Prediction\nUsing Machine Learning",
      0.6, 1.1, 9.8, 2.1, size=40, bold=True, color=WHITE)
    T(sl, "A Medical Decision Support System",
      0.6, 3.15, 9.0, 0.55, size=18, italic=True, color=LIGHT_BLUE)

    R(sl, 0.6, 4.55, 4.0, 0.055, ORANGE)

    T(sl, "Jamshaid Amjad",
      0.6, 4.7, 5.5, 0.5, size=17, bold=True, color=WHITE)
    T(sl, "AI in Healthcare — Capstone Project",
      0.6, 5.22, 6.5, 0.38, size=13, color=LIGHT_BLUE)
    T(sl, "May 2026",
      0.6, 5.62, 4.0, 0.38, size=13, color=LIGHT_BLUE)

    # Bottom stats strip
    R(sl, 0, 6.6, 13.33, 0.9, RGBColor(0x14, 0x24, 0x42))
    for i, (val, lbl) in enumerate([
        ("768", "Patients"), ("8", "Clinical Features"),
        ("3", "ML Models"), ("87.5%", "Best ROC-AUC")
    ]):
        x = 1.2 + i * 2.8
        T(sl, val, x, 6.65, 2.4, 0.38, size=20, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
        T(sl, lbl, x, 7.04, 2.4, 0.35, size=10, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Agenda
# ═══════════════════════════════════════════════════════════════════════════════
def s02_agenda(prs):
    sl = _blank(prs)
    R(sl, 0, 0, 13.33, 7.5, PALE_BLUE)
    header(sl, "Agenda", "What this presentation covers — 15 slides")
    footer(sl, 2)

    sections = [
        ("01", "Problem & Clinical Motivation",  "Why diabetes screening matters at scale",          "3"),
        ("02", "Dataset & Exploratory Analysis",  "Pima Indians dataset, EDA, class imbalance",      "4–6"),
        ("03", "Preprocessing Pipeline",          "6-step leakage-free data preparation",            "7"),
        ("04", "Feature Selection",               "ANOVA + Random Forest combined ranking",          "8"),
        ("05", "Model Architectures",             "Logistic Regression, Random Forest, MLP",         "9–11"),
        ("06", "Training Strategy",               "Cross-validation, GridSearchCV, early stopping",  "12"),
        ("07", "Results & Evaluation",            "Full metrics, ROC curves, model comparison",      "13–14"),
        ("08", "Conclusion & Future Work",        "Key findings, limitations, clinical relevance",   "15"),
    ]

    col_w = 5.9
    for i, (num, title, desc, slides) in enumerate(sections):
        row = i // 2; col = i % 2
        l = 0.35 + col * (col_w + 0.53)
        t = 1.25 + row * 1.4
        R(sl, l, t, col_w, 1.25, WHITE, line_color=RULE_GRAY, line_w=0.5)
        R(sl, l, t, 0.55, 1.25, NAVY)
        T(sl, num, l+0.04, t+0.38, 0.48, 0.52, size=16, bold=True,
          color=ORANGE, align=PP_ALIGN.CENTER)
        T(sl, title, l+0.65, t+0.10, col_w-1.1, 0.42, size=13, bold=True, color=NAVY)
        T(sl, desc,  l+0.65, t+0.52, col_w-1.1, 0.38, size=11, color=MED_GRAY)
        R(sl, l+col_w-0.72, t+0.85, 0.62, 0.30, LIGHT_BLUE)
        T(sl, f"Slide {slides}", l+col_w-0.70, t+0.88, 0.62, 0.24,
          size=9, color=BLUE, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Problem & Clinical Motivation
# ═══════════════════════════════════════════════════════════════════════════════
def s03_problem(prs):
    sl = _blank(prs)
    R(sl, 0, 0, 13.33, 7.5, PALE_BLUE)
    header(sl, "Problem & Clinical Motivation",
           "Diabetes is a global epidemic — early detection saves lives")
    footer(sl, 3)

    # Stat cards
    for i, (v, l) in enumerate([
        ("537M",   "Adults with diabetes\nglobally (2021)"),
        ("783M",   "Projected by 2045\n(+46% increase)"),
        (">90%",   "Are Type 2 diabetes\n(preventable)"),
        ("$966B",  "Annual global health\nexpenditure on diabetes"),
    ]):
        stat_card(sl, v, l, 0.35 + i*3.17, 1.22, w=2.9, h=1.5)

    # Two column content
    bullet_card(sl, "The Clinical Challenge", [
        "Type 2 diabetes progresses silently for years",
        "Often diagnosed after irreversible organ damage",
        "Standard fasting glucose test is invasive & costly",
        "Primary care is overwhelmed — needs decision tools",
        "Timely screening reduces complications by up to 58%",
    ], 0.35, 2.85, 6.0, 3.1, stripe=ORANGE)

    bullet_card(sl, "How Machine Learning Can Help", [
        "Predict risk from routine clinical blood measurements",
        "Automate flag of high-risk patients for confirmatory test",
        "Provide probabilistic scores, not just binary outputs",
        "Interpretable models support clinical decision-making",
        "Cost-effective: uses data already collected at primary care",
    ], 6.63, 2.85, 6.35, 3.1, stripe=GREEN)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Dataset Overview
# ═══════════════════════════════════════════════════════════════════════════════
def s04_dataset(prs):
    sl = _blank(prs)
    R(sl, 0, 0, 13.33, 7.5, PALE_BLUE)
    header(sl, "Dataset Overview",
           "Pima Indians Diabetes Dataset — NIDDK / UCI Machine Learning Repository")
    footer(sl, 4)

    # Left: dataset meta cards
    for i, (v, l) in enumerate([
        ("768",  "Total patients"),
        ("8",    "Input features"),
        ("268",  "Positive cases (35%)"),
        ("500",  "Negative cases (65%)"),
    ]):
        metric_highlight(sl, l, v, 0.35, 1.25+i*1.35, w=3.2, h=1.22,
                         stripe=BLUE if i%2==0 else ORANGE)

    # Right: features table
    R(sl, 3.85, 1.25, 9.13, 5.75, WHITE, line_color=RULE_GRAY, line_w=0.5)
    T(sl, "Clinical Features", 3.98, 1.30, 8.85, 0.42, size=13, bold=True, color=NAVY)

    # Table header
    col_positions = [3.85, 6.55, 9.25, 11.25]
    col_widths    = [2.68, 2.68, 2.0,  1.73]
    col_headers   = ["Feature", "Description", "Unit / Type", "Missing?"]
    R(sl, 3.85, 1.75, 9.13, 0.42, NAVY)
    for j, (cx, cw, ch) in enumerate(zip(col_positions, col_widths, col_headers)):
        T(sl, ch, cx+0.06, 1.77, cw-0.1, 0.38,
          size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    rows = [
        ("Glucose",                  "Plasma glucose conc. (2hr OGTT)", "mg/dL",  "Yes (4.9%)"),
        ("BMI",                      "Body Mass Index",                  "kg/m²",  "Yes (1.4%)"),
        ("BloodPressure",            "Diastolic blood pressure",         "mm Hg",  "Yes (4.6%)"),
        ("Insulin",                  "2-hour serum insulin",             "μU/mL",  "Yes (48.7%)"),
        ("SkinThickness",            "Triceps skin-fold thickness",      "mm",     "Yes (29.6%)"),
        ("DiabetesPedigreeFunction", "Genetic diabetes risk score",      "0–2.42", "No"),
        ("Age",                      "Age of patient",                   "Years",  "No"),
        ("Pregnancies",              "Number of pregnancies",            "Count",  "No"),
        ("Outcome  (TARGET)",        "Diabetes diagnosis",               "0 / 1",  "—"),
    ]
    for ri, row in enumerate(rows):
        bg = LIGHT_BLUE if ri % 2 == 0 else WHITE
        t_row = 2.18 + ri * 0.51
        R(sl, 3.85, t_row, 9.13, 0.51, bg)
        for j, (cx, cw, val) in enumerate(zip(col_positions, col_widths, row)):
            bold = (ri == 8)
            color = NAVY if ri == 8 else DARK_GRAY
            T(sl, val, cx+0.06, t_row+0.08, cw-0.1, 0.38,
              size=10, bold=bold, color=color,
              align=PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — EDA: Class Distribution & Missing Data
# ═══════════════════════════════════════════════════════════════════════════════
def s05_eda1(prs):
    sl = _blank(prs)
    R(sl, 0, 0, 13.33, 7.5, PALE_BLUE)
    header(sl, "Exploratory Data Analysis — Class Distribution & Missing Values",
           "Understanding imbalance and data quality before modelling")
    footer(sl, 5)

    # Left: class distribution figure
    card(sl, 0.35, 1.22, 5.7, 5.68)
    T(sl, "Class Distribution", 0.49, 1.30, 5.42, 0.42, size=13, bold=True, color=NAVY)
    PIC(sl, "class_distribution.png", 0.5, 1.78, w=5.4)

    # Right: missing data analysis
    card(sl, 6.28, 1.22, 6.7, 2.68, stripe=ORANGE)
    T(sl, "Missing Data (Biological Zeros → NaN)", 6.42, 1.30, 6.42, 0.42,
      size=13, bold=True, color=NAVY)
    missing = [
        ("Insulin",        "48.7%", RED),
        ("SkinThickness",  "29.6%", ORANGE),
        ("BloodPressure",  "4.6%",  BLUE),
        ("Glucose",        "4.9%",  BLUE),
        ("BMI",            "1.4%",  GREEN),
    ]
    for i, (feat, pct, bar_col) in enumerate(missing):
        ty = 1.78 + i * 0.46
        T(sl, feat, 6.42, ty, 2.3, 0.38, size=11, color=DARK_GRAY)
        bar_w = float(pct.replace("%","")) / 100 * 3.4
        R(sl, 8.85, ty+0.06, 3.4, 0.25, LIGHT_GRAY)
        R(sl, 8.85, ty+0.06, bar_w, 0.25, bar_col)
        T(sl, pct, 12.35, ty, 0.55, 0.38, size=11, bold=True, color=bar_col)

    # Right bottom: key insights
    card(sl, 6.28, 4.05, 6.7, 2.85, stripe=GREEN)
    T(sl, "Key EDA Insights", 6.42, 4.13, 6.42, 0.42, size=13, bold=True, color=NAVY)
    insights = [
        "65% negative / 35% positive → class imbalance requires SMOTE",
        "Insulin has the most missing data (49%) — median imputation used",
        "Zeros in Glucose / BMI are biologically impossible → treated as NaN",
        "All features are continuous except Pregnancies & Outcome (discrete)",
        "No duplicate rows; dataset is complete after zero replacement",
    ]
    for i, ins in enumerate(insights):
        T(sl, f"▸  {ins}", 6.42, 4.60+i*0.44, 6.42, 0.40, size=11, color=DARK_GRAY)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — EDA: Feature Distributions & Correlations
# ═══════════════════════════════════════════════════════════════════════════════
def s06_eda2(prs):
    sl = _blank(prs)
    R(sl, 0, 0, 13.33, 7.5, PALE_BLUE)
    header(sl, "Exploratory Data Analysis — Feature Distributions & Correlations",
           "Understanding which features separate diabetic from non-diabetic patients")
    footer(sl, 6)

    card(sl, 0.35, 1.22, 7.7, 5.68)
    T(sl, "Feature Distributions by Outcome (KDE)", 0.49, 1.30, 7.42, 0.42,
      size=13, bold=True, color=NAVY)
    PIC(sl, "feature_distributions.png", 0.42, 1.78, w=7.55)

    card(sl, 8.28, 1.22, 4.7, 5.68, stripe=ORANGE)
    T(sl, "Correlation with Outcome", 8.42, 1.30, 4.42, 0.42, size=13, bold=True, color=NAVY)

    corr_data = [
        ("Glucose",                  0.47, BLUE),
        ("BMI",                      0.29, BLUE),
        ("Age",                      0.24, BLUE),
        ("Pregnancies",              0.22, BLUE),
        ("DiabetesPedigreeFunction", 0.17, BLUE),
        ("Insulin",                  0.13, BLUE),
        ("SkinThickness",            0.07, MED_GRAY),
        ("BloodPressure",            0.07, MED_GRAY),
    ]
    for i, (feat, corr, bar_col) in enumerate(corr_data):
        ty = 1.78 + i * 0.56
        T(sl, feat, 8.42, ty, 2.45, 0.38, size=10, color=DARK_GRAY)
        bar_w = corr / 0.5 * 1.85
        R(sl, 10.96, ty+0.06, 1.85, 0.25, LIGHT_GRAY)
        R(sl, 10.96, ty+0.06, bar_w, 0.25, bar_col)
        T(sl, f"{corr:.2f}", 12.88, ty, 0.5, 0.38, size=10, bold=True, color=bar_col)

    T(sl, "Glucose is the strongest predictor (r=0.47),\nfollowed by BMI and Age.",
      8.42, 6.52, 4.42, 0.55, size=10, italic=True, color=MED_GRAY)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Preprocessing Pipeline
# ═══════════════════════════════════════════════════════════════════════════════
def s07_pipeline(prs):
    sl = _blank(prs)
    R(sl, 0, 0, 13.33, 7.5, PALE_BLUE)
    header(sl, "Preprocessing Pipeline",
           "Six leakage-free steps — all transformers fitted on training data only")
    footer(sl, 7)

    steps = [
        ("01", "Load &\nClean",
         "Replace biological zeros\n(Glucose, BMI, Insulin,\nBloodPressure, SkinThickness)\nwith NaN"),
        ("02", "Stratified\nSplit",
         "70% Train\n15% Validation\n15% Test\nStratified by Outcome"),
        ("03", "Median\nImputation",
         "SimpleImputer (median)\nFit on train ONLY\nApplied to val + test\nNo leakage"),
        ("04", "StandardScaler\nNormalisation",
         "Mean=0, Std=1\nFit on train ONLY\nApplied to val + test\nEssential for LR & MLP"),
        ("05", "Feature\nSelection",
         "Top-5 by average of:\n• ANOVA F-test rank\n• Random Forest\n  importance rank"),
        ("06", "SMOTE\nOversampling",
         "Synthetic Minority\nOversampling — train\nset ONLY\nBalances 65/35 → 50/50"),
    ]

    box_w = 1.95
    gap   = 0.22
    start = 0.35
    for i, (num, title, desc) in enumerate(steps):
        x = start + i * (box_w + gap)
        step_box(sl, num, title, desc, x, 1.22, w=box_w, h=2.55)
        if i < 5:
            T(sl, "▶", x + box_w + 0.02, 1.22 + 1.0, 0.22, 0.55,
              size=16, color=BLUE, align=PP_ALIGN.CENTER)

    # Bottom: why leakage-free matters
    R(sl, 0.35, 3.92, 12.63, 2.98, WHITE, line_color=RULE_GRAY, line_w=0.5)
    R(sl, 0.35, 3.92, 12.63, 0.055, ORANGE)
    T(sl, "Why Leakage-Free Matters", 0.5, 3.99, 6.0, 0.42,
      size=13, bold=True, color=NAVY)

    left_points = [
        "✗  Fitting scaler on ALL data leaks test statistics into training",
        "✗  Imputing on ALL data uses future test values to fill training gaps",
        "✗  SMOTE applied to test set creates artificial, non-real samples",
    ]
    right_points = [
        "✓  All transformers fit on train set — same transform applied to val/test",
        "✓  Stratified split preserves class ratio in every split",
        "✓  This pipeline exactly mirrors real deployment conditions",
    ]
    for i, pt in enumerate(left_points):
        T(sl, pt, 0.5, 4.52+i*0.72, 6.1, 0.62, size=11, color=DARK_GRAY)
    for i, pt in enumerate(right_points):
        T(sl, pt, 6.8, 4.52+i*0.72, 6.0, 0.62, size=11, color=DARK_GRAY)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Feature Selection
# ═══════════════════════════════════════════════════════════════════════════════
def s08_features(prs):
    sl = _blank(prs)
    R(sl, 0, 0, 13.33, 7.5, PALE_BLUE)
    header(sl, "Feature Selection",
           "Combining ANOVA F-test and Random Forest importance for robust top-5 selection")
    footer(sl, 8)

    # Left: methodology
    card(sl, 0.35, 1.22, 5.5, 5.68, stripe=ORANGE)
    T(sl, "Selection Methodology", 0.49, 1.30, 5.22, 0.42, size=13, bold=True, color=NAVY)

    method_steps = [
        ("Step 1", "ANOVA F-test",
         "Compute F-statistic for each feature\nvs. Outcome. Rank 1–8 (1=highest)."),
        ("Step 2", "Random Forest Rank",
         "Train RF on training data. Use\nfeature_importances_ to rank 1–8."),
        ("Step 3", "Average Rank",
         "Average the two ranks per feature.\nSelect the 5 features with lowest\n(best) average rank."),
        ("Step 4", "Apply Selection",
         "Use the same 5 feature indices on\nval + test sets. Never refit on test."),
    ]
    for i, (step, title, desc) in enumerate(method_steps):
        t_y = 1.78 + i * 1.22
        R(sl, 0.49, t_y, 1.0, 0.98, NAVY)
        T(sl, step, 0.52, t_y+0.08, 0.95, 0.35, size=9, bold=True,
          color=ORANGE, align=PP_ALIGN.CENTER)
        T(sl, title, 1.58, t_y+0.05, 3.9, 0.32, size=11, bold=True, color=NAVY)
        T(sl, desc,  1.58, t_y+0.36, 3.9, 0.62, size=10, color=DARK_GRAY)

    # Right: feature importance chart
    card(sl, 6.1, 1.22, 6.88, 5.68, stripe=BLUE)
    T(sl, "Feature Importances (Random Forest)", 6.24, 1.30, 6.6, 0.42,
      size=13, bold=True, color=NAVY)
    PIC(sl, "feature_importance.png", 6.18, 1.78, w=6.72)

    # Bottom callout
    R(sl, 6.1, 6.2, 6.88, 0.62, NAVY)
    T(sl, "Selected Features: Glucose · BMI · Insulin · Age · Pregnancies",
      6.24, 6.27, 6.6, 0.48, size=12, bold=True,
      color=ORANGE, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Model 1: Logistic Regression
# ═══════════════════════════════════════════════════════════════════════════════
def s09_lr(prs):
    sl = _blank(prs)
    R(sl, 0, 0, 13.33, 7.5, PALE_BLUE)
    header(sl, "Model 1 — Logistic Regression",
           "Interpretable linear baseline — essential reference point for clinical deployment")
    footer(sl, 9)

    # Left: description
    bullet_card(sl, "Model Description", [
        "Linear model: log-odds = β₀ + β₁x₁ + … + β₅x₅",
        "Solver: L-BFGS (limited-memory quasi-Newton)",
        "Regularisation: L2 penalty (C=1.0 default)",
        "Max iterations: 1000 (ensures convergence)",
        "Threshold: 0.5 for binary prediction",
    ], 0.35, 1.22, 6.0, 3.05, stripe=BLUE)

    bullet_card(sl, "Why Logistic Regression?", [
        "Fully interpretable — coefficients show feature direction",
        "Clinicians can understand and trust the prediction",
        "Fast to train and deploy in resource-limited settings",
        "Calibrated probabilities useful for risk stratification",
        "Strong baseline: if complex models barely beat it, use LR",
    ], 0.35, 4.42, 6.0, 2.48, stripe=ORANGE)

    # Right: results
    card(sl, 6.63, 1.22, 6.35, 5.68, stripe=BLUE)
    T(sl, "Test-Set Performance (n=116)", 6.77, 1.30, 6.07, 0.42,
      size=13, bold=True, color=NAVY)

    metrics = [
        ("Accuracy",   "77.6%", BLUE),
        ("Precision",  "68.3%", BLUE),
        ("Recall",     "68.3%", BLUE),
        ("F1 Score",   "68.3%", BLUE),
        ("ROC-AUC",    "0.866", ORANGE),
    ]
    for i, (lbl, val, col) in enumerate(metrics):
        row = i // 2; col_n = i % 2
        l = 6.77 + col_n * 3.15
        t = 1.82 + row * 1.12
        w = 2.9
        metric_highlight(sl, lbl, val, l, t, w=w, h=0.95, val_color=col, stripe=col)

    # CV result
    R(sl, 6.63, 5.1, 6.35, 0.8, NAVY)
    T(sl, "5-Fold Cross-Validation F1",
      6.77, 5.15, 6.07, 0.32, size=11, bold=True, color=LIGHT_BLUE)
    T(sl, "0.718  ±  0.029",
      6.77, 5.45, 6.07, 0.38, size=18, bold=True,
      color=ORANGE, align=PP_ALIGN.CENTER)

    # Confusion matrix
    T(sl, "Confusion Matrix", 6.77, 6.02, 6.07, 0.38, size=11, bold=True, color=NAVY)
    PIC(sl, "logistic_regression_confusion.png", 7.2, 6.02, w=5.6)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Model 2: Random Forest
# ═══════════════════════════════════════════════════════════════════════════════
def s10_rf(prs):
    sl = _blank(prs)
    R(sl, 0, 0, 13.33, 7.5, PALE_BLUE)
    header(sl, "Model 2 — Random Forest",
           "Ensemble of decision trees with GridSearchCV hyperparameter optimisation")
    footer(sl, 10)

    bullet_card(sl, "Model Description", [
        "Ensemble of N independent decision trees (bagging)",
        "Each tree trained on bootstrap sample of training data",
        "Each split considers random subset of features (√5 ≈ 2)",
        "Final prediction: majority vote across all trees",
        "Naturally handles non-linear interactions between features",
    ], 0.35, 1.22, 6.0, 3.0, stripe=GREEN)

    # GridSearch results box
    card(sl, 0.35, 4.37, 6.0, 2.53, stripe=ORANGE)
    T(sl, "GridSearchCV Results (5-fold)", 0.49, 4.45, 5.72, 0.42,
      size=13, bold=True, color=NAVY)
    grid_rows = [
        ("Parameter",        "Search Space",        "Best Value"),
        ("n_estimators",     "100, 300",             "100"),
        ("max_depth",        "None, 10, 20",         "10"),
        ("min_samples_split","2, 5",                 "2"),
        ("min_samples_leaf", "1, 2",                 "1"),
    ]
    for ri, row in enumerate(grid_rows):
        t_y = 4.93 + ri * 0.38
        bg = NAVY if ri == 0 else (LIGHT_BLUE if ri % 2 == 1 else WHITE)
        R(sl, 0.35, t_y, 6.0, 0.38, bg)
        for ci, (val, cw, cx) in enumerate(zip(row, [2.1, 2.1, 1.7], [0.35, 2.47, 4.59])):
            fc = WHITE if ri == 0 else DARK_GRAY
            bold = (ri == 0)
            T(sl, val, cx+0.07, t_y+0.06, cw-0.12, 0.28,
              size=10, bold=bold, color=fc)

    # Right panel
    card(sl, 6.63, 1.22, 6.35, 5.68, stripe=GREEN)
    T(sl, "Test-Set Performance (n=116)", 6.77, 1.30, 6.07, 0.42,
      size=13, bold=True, color=NAVY)

    metrics = [
        ("Accuracy",  "76.7%", GREEN),
        ("Precision", "69.4%", GREEN),
        ("Recall",    "61.0%", GREEN),
        ("F1 Score",  "64.9%", GREEN),
        ("ROC-AUC",   "0.841", ORANGE),
    ]
    for i, (lbl, val, col) in enumerate(metrics):
        row = i // 2; col_n = i % 2
        l = 6.77 + col_n * 3.15
        t = 1.82 + row * 1.12
        metric_highlight(sl, lbl, val, l, t, w=2.9, h=0.95, val_color=col, stripe=col)

    R(sl, 6.63, 5.1, 6.35, 0.8, NAVY)
    T(sl, "5-Fold Cross-Validation F1",
      6.77, 5.15, 6.07, 0.32, size=11, bold=True, color=LIGHT_BLUE)
    T(sl, "0.801  ±  0.042",
      6.77, 5.45, 6.07, 0.38, size=18, bold=True,
      color=ORANGE, align=PP_ALIGN.CENTER)

    T(sl, "Confusion Matrix", 6.77, 6.02, 6.07, 0.38, size=11, bold=True, color=NAVY)
    PIC(sl, "random_forest_confusion.png", 7.2, 6.02, w=5.6)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Model 3: MLP Neural Network
# ═══════════════════════════════════════════════════════════════════════════════
def s11_mlp(prs):
    sl = _blank(prs)
    R(sl, 0, 0, 13.33, 7.5, PALE_BLUE)
    header(sl, "Model 3 — MLP Neural Network",
           "Feed-forward network with early stopping and learning-rate scheduling")
    footer(sl, 11)

    # Architecture visual (left)
    card(sl, 0.35, 1.22, 6.0, 5.68, stripe=ORANGE)
    T(sl, "Network Architecture", 0.49, 1.30, 5.72, 0.42, size=13, bold=True, color=NAVY)

    arch_layers = [
        ("Input Layer",      "5 neurons (selected features)",  BLUE,    1.88),
        ("Dense Layer 1",    "16 neurons + ReLU activation",   ORANGE,  2.60),
        ("Dropout",          "rate = 0.2  (regularisation)",   MED_GRAY,3.18),
        ("Dense Layer 2",    "8 neurons + ReLU activation",    ORANGE,  3.76),
        ("Output Layer",     "1 neuron + Sigmoid → P(diabetes)",GREEN,  4.44),
    ]
    for lbl, desc, col, ty in arch_layers:
        R(sl, 0.55, ty, 5.6, 0.52, col)
        T(sl, lbl,  0.65, ty+0.05, 2.2, 0.42, size=11, bold=True, color=WHITE)
        T(sl, desc, 2.88, ty+0.05, 3.2, 0.42, size=10, color=WHITE)
        if ty < 4.44:
            T(sl, "↓", 2.9, ty+0.52, 0.6, 0.3, size=14, color=NAVY, align=PP_ALIGN.CENTER)

    # Training config
    train_items = [
        ("Loss function",    "BCEWithLogitsLoss (numerically stable)"),
        ("Optimiser",        "Adam  (lr=0.001)"),
        ("LR Scheduler",     "ReduceLROnPlateau (factor=0.5, patience=5)"),
        ("Epochs",           "Up to 50 (early stopping at patience=10)"),
        ("Batch size",       "32"),
    ]
    for i, (k, v) in enumerate(train_items):
        ty = 5.0 + i * 0.38
        T(sl, f"{k}:", 0.49, ty, 2.2, 0.35, size=10, bold=True, color=NAVY)
        T(sl, v,       2.72, ty, 3.5, 0.35, size=10, color=DARK_GRAY)

    # Right: results + training curve
    card(sl, 6.63, 1.22, 6.35, 5.68, stripe=ORANGE)
    T(sl, "Test-Set Performance (n=116)", 6.77, 1.30, 6.07, 0.42,
      size=13, bold=True, color=NAVY)

    metrics = [
        ("Accuracy",  "76.7%", ORANGE),
        ("Precision", "67.5%", ORANGE),
        ("Recall",    "65.9%", ORANGE),
        ("F1 Score",  "66.7%", ORANGE),
        ("ROC-AUC",   "0.875 ★", GREEN),
    ]
    for i, (lbl, val, col) in enumerate(metrics):
        row = i // 2; col_n = i % 2
        l = 6.77 + col_n * 3.15
        t = 1.82 + row * 1.12
        metric_highlight(sl, lbl, val, l, t, w=2.9, h=0.95, val_color=col, stripe=col)

    R(sl, 6.63, 5.1, 6.35, 0.52, NAVY)
    T(sl, "★  Highest ROC-AUC of all three models",
      6.77, 5.14, 6.07, 0.42, size=11, bold=True, color=ORANGE)

    T(sl, "Training Curves", 6.77, 5.74, 6.07, 0.32, size=11, bold=True, color=NAVY)
    PIC(sl, "mlp_training.png", 6.68, 5.74, w=6.27)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Training Strategy & Cross-Validation
# ═══════════════════════════════════════════════════════════════════════════════
def s12_training(prs):
    sl = _blank(prs)
    R(sl, 0, 0, 13.33, 7.5, PALE_BLUE)
    header(sl, "Training Strategy & Cross-Validation",
           "Rigorous evaluation methodology to ensure reliable, unbiased results")
    footer(sl, 12)

    # Eval strategy overview
    card(sl, 0.35, 1.22, 12.63, 1.38, stripe=ORANGE)
    T(sl, "Evaluation Strategy", 0.49, 1.30, 12.35, 0.42, size=13, bold=True, color=NAVY)
    for i, (title, desc) in enumerate([
        ("Held-out Test Set",  "15% of data (n=116) — never seen during training or tuning"),
        ("Validation Set",     "15% for early stopping & hyperparameter guidance (n=115)"),
        ("Train Set",          "70% for fitting (n=538 raw, n=700 after SMOTE)"),
    ]):
        T(sl, f"▸  {title}:", 0.49+i*4.25, 1.72, 2.0, 0.32, size=10, bold=True, color=NAVY)
        T(sl, desc, 0.49+i*4.25+2.0, 1.72, 2.1, 0.58, size=10, color=DARK_GRAY)

    # Left: cross-validation
    card(sl, 0.35, 2.75, 6.1, 4.15, stripe=BLUE)
    T(sl, "5-Fold Cross-Validation (LR & RF)", 0.49, 2.83, 5.82, 0.42,
      size=13, bold=True, color=NAVY)
    cv_steps = [
        "1.  Combine train + validation set (n=653)",
        "2.  Split into 5 equal folds (~131 samples each)",
        "3.  Train on 4 folds, evaluate on 1 fold — repeat 5×",
        "4.  Report mean F1 ± standard deviation",
        "5.  Best RF params from GridSearch applied to final model",
    ]
    for i, s in enumerate(cv_steps):
        T(sl, s, 0.49, 3.35+i*0.52, 5.82, 0.48, size=11, color=DARK_GRAY)

    # CV Results table
    R(sl, 0.35, 5.67, 6.1, 0.36, NAVY)
    for ci, (h, w, x) in enumerate(zip(["Model","CV F1 Mean","Std Dev","Fold Scores"],
                                        [2.1, 1.3, 1.0, 1.5],
                                        [0.35, 2.47, 3.79, 4.81])):
        T(sl, h, x+0.07, 5.70, w-0.1, 0.28, size=10, bold=True, color=WHITE)

    for ri, (name, mean, std, folds) in enumerate([
        ("Logistic Regression", "0.718", "±0.029", "[0.736, 0.685, 0.682, 0.748, 0.739]"),
        ("Random Forest",       "0.801", "±0.042", "[0.756, 0.769, 0.840, 0.862, 0.778]"),
    ]):
        bg = LIGHT_BLUE if ri == 0 else WHITE
        R(sl, 0.35, 6.04+ri*0.58, 6.1, 0.55, bg)
        for ci, (val, w, x) in enumerate(zip([name, mean, std, folds],
                                              [2.1, 1.3, 1.0, 1.5],
                                              [0.35, 2.47, 3.79, 4.81])):
            T(sl, val, x+0.07, 6.07+ri*0.58, w-0.1, 0.42,
              size=10, color=DARK_GRAY, bold=(ci==0))

    # Right: GridSearch + SMOTE
    card(sl, 6.73, 2.75, 6.25, 2.15, stripe=GREEN)
    T(sl, "GridSearchCV — Random Forest", 6.87, 2.83, 5.97, 0.42,
      size=13, bold=True, color=NAVY)
    gs_rows = [("Search space", "24 combinations (2×3×2×2)"),
               ("CV folds",     "5-fold (same as CV above)"),
               ("Scoring",      "F1 score"),
               ("Best CV F1",   "0.8299"),
               ("Best params",  "max_depth=10, n_estimators=100")]
    for i, (k, v) in enumerate(gs_rows):
        T(sl, f"{k}:", 6.87, 3.33+i*0.38, 2.4, 0.35, size=10, bold=True, color=NAVY)
        T(sl, v,       9.3,  3.33+i*0.38, 3.5, 0.35, size=10, color=DARK_GRAY)

    card(sl, 6.73, 5.05, 6.25, 1.85, stripe=ORANGE)
    T(sl, "SMOTE — Class Balancing", 6.87, 5.13, 5.97, 0.42, size=13, bold=True, color=NAVY)
    TM(sl, [
        "▸  Training set: 65% negative / 35% positive (imbalanced)",
        "▸  SMOTE generates synthetic minority (diabetic) samples",
        "▸  After SMOTE: 350 positive / 350 negative — 50/50",
        "▸  Applied to training set ONLY — val/test remain real data",
    ], 6.87, 5.57, 5.97, 1.25, size=11, color=DARK_GRAY, spacing=4)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Performance Results
# ═══════════════════════════════════════════════════════════════════════════════
def s13_results(prs):
    sl = _blank(prs)
    R(sl, 0, 0, 13.33, 7.5, PALE_BLUE)
    header(sl, "Performance Results — Test Set (n=116)",
           "Held-out evaluation on 15% of data — never seen during training or hyperparameter selection")
    footer(sl, 13)

    # Main results table
    R(sl, 0.35, 1.22, 12.63, 3.42, WHITE, line_color=RULE_GRAY, line_w=0.5)
    T(sl, "Complete Metrics — All Models vs. Majority Baseline",
      0.49, 1.27, 12.35, 0.42, size=13, bold=True, color=NAVY)

    col_x = [0.35, 3.2,  5.35, 7.25, 9.1, 10.85, 12.3]
    col_w = [2.82, 2.12, 1.88, 1.82, 1.72,  1.42,  0.68]
    headers = ["Model", "Accuracy", "Precision", "Recall", "F1 Score", "ROC-AUC", "Δ Base"]
    R(sl, 0.35, 1.72, 12.63, 0.46, NAVY)
    for cx, cw, h in zip(col_x, col_w, headers):
        T(sl, h, cx+0.08, 1.75, cw-0.14, 0.38,
          size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    rows = [
        ("Logistic Regression", "0.776", "0.683", "0.683", "0.683", "0.866", "+36.6%", BLUE),
        ("Random Forest",       "0.767", "0.694", "0.610", "0.649", "0.841", "+34.1%", GREEN),
        ("MLP Neural Network",  "0.767", "0.675", "0.659", "0.667", "0.875", "+37.5%", ORANGE),
        ("Majority Baseline",   "0.647", "N/A",   "0.000", "0.000", "0.500", "—",      MED_GRAY),
    ]
    for ri, (name, acc, prec, rec, f1, auc, delta, col) in enumerate(rows):
        bg = RGBColor(0xF8, 0xFB, 0xFF) if ri % 2 == 0 else WHITE
        R(sl, 0.35, 2.19+ri*0.44, 12.63, 0.44, bg)
        vals = [name, acc, prec, rec, f1, auc, delta]
        is_best = [False,
                   ri == 0,   # LR best acc
                   ri == 1,   # RF best precision
                   ri == 0,   # LR best recall
                   ri == 0,   # LR best F1
                   ri == 2,   # MLP best AUC
                   False]
        for ci, (cx, cw, val, best) in enumerate(zip(col_x, col_w, vals, is_best)):
            cell_col = col if ci == 0 else (GREEN if best else DARK_GRAY)
            bold = ci == 0 or best
            T(sl, val, cx+0.08, 2.22+ri*0.44, cw-0.14, 0.36,
              size=11 if ci > 0 else 11, bold=bold, color=cell_col,
              align=PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER)

    T(sl, "★ Green values = best in column   |   Δ Base = improvement over majority-class baseline AUC",
      0.49, 3.90, 12.35, 0.38, size=9, italic=True, color=MED_GRAY)

    # Bottom: 4 insight callouts
    highlights = [
        (BLUE,   "Best Overall Accuracy", "Logistic Regression\n77.6%"),
        (GREEN,  "Best Precision",        "Random Forest\n69.4%"),
        (ORANGE, "Best ROC-AUC",          "MLP Neural Net\n0.875"),
        (NAVY,   "Best CV Stability",     "Random Forest\nF1 0.801 ± 0.042"),
    ]
    for i, (col, lbl, val) in enumerate(highlights):
        l = 0.35 + i * 3.22
        R(sl, l, 4.35, 3.0, 2.55, col)
        R(sl, l, 4.35, 3.0, 0.055, ORANGE)
        T(sl, lbl, l+0.14, 4.45, 2.72, 0.48, size=11, bold=True, color=WHITE)
        T(sl, val, l+0.14, 5.05, 2.72, 1.0,  size=22, bold=True,
          color=WHITE, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Model Comparison Visualisations
# ═══════════════════════════════════════════════════════════════════════════════
def s14_comparison(prs):
    sl = _blank(prs)
    R(sl, 0, 0, 13.33, 7.5, PALE_BLUE)
    header(sl, "Model Comparison — ROC Curves & Performance Charts",
           "Visual evaluation of discrimination ability and metric trade-offs")
    footer(sl, 14)

    # ROC curves (left)
    card(sl, 0.35, 1.22, 6.6, 5.68, stripe=BLUE)
    T(sl, "ROC Curves — All Models vs. Baseline", 0.49, 1.30, 6.32, 0.42,
      size=13, bold=True, color=NAVY)
    PIC(sl, "roc_curves_comparison.png", 0.42, 1.78, w=6.46)

    # AUC summary in the card footer
    R(sl, 0.35, 6.28, 6.6, 0.52, NAVY)
    for i, (m, auc, col) in enumerate([
        ("LR", "0.866", BLUE), ("RF", "0.841", GREEN), ("MLP", "0.875", ORANGE)
    ]):
        T(sl, f"{m}: AUC={auc}", 0.55+i*2.17, 6.33, 2.1, 0.38,
          size=12, bold=True, color=col, align=PP_ALIGN.CENTER)

    # Bar chart (right)
    card(sl, 7.22, 1.22, 5.76, 3.68, stripe=GREEN)
    T(sl, "Accuracy / F1 / ROC-AUC by Model", 7.36, 1.30, 5.48, 0.42,
      size=13, bold=True, color=NAVY)
    PIC(sl, "model_comparison.png", 7.28, 1.78, w=5.62)

    # Interpretation box
    card(sl, 7.22, 5.05, 5.76, 1.85, stripe=ORANGE)
    T(sl, "Interpretation", 7.36, 5.13, 5.48, 0.42, size=13, bold=True, color=NAVY)
    TM(sl, [
        "▸  All models substantially outperform random (AUC=0.500)",
        "▸  MLP leads on ROC-AUC — best probabilistic triage tool",
        "▸  LR leads on F1 — best balanced classifier for deployment",
        "▸  RF most stable across CV folds (std=0.042)",
        "▸  Accuracy differences small (~1%) — ROC-AUC is the key metric",
    ], 7.36, 5.57, 5.48, 1.25, size=11, color=DARK_GRAY, spacing=3)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — Conclusion & Future Work
# ═══════════════════════════════════════════════════════════════════════════════
def s15_conclusion(prs):
    sl = _blank(prs)
    R(sl, 0, 0, 13.33, 7.5, NAVY)
    R(sl, 0, 6.55, 13.33, 0.95, RGBColor(0x14, 0x24, 0x42))
    R(sl, 0, 0, 0.07, 7.5, ORANGE)

    # Title band
    R(sl, 0.07, 0, 13.26, 1.12, RGBColor(0x17, 0x2B, 0x52))
    T(sl, "Conclusion & Future Work", 0.25, 0.08, 12.8, 0.66,
      size=27, bold=True, color=WHITE)
    T(sl, "Key findings, clinical implications, limitations, and next steps",
      0.25, 0.74, 12.8, 0.34, size=12, italic=True, color=LIGHT_BLUE)

    # 4 Key findings (top)
    T(sl, "Key Findings", 0.25, 1.22, 5.5, 0.42, size=14, bold=True, color=ORANGE)
    findings = [
        ("All models ROC-AUC > 0.84",
         "Routine clinical measurements are sufficient for meaningful diabetes screening"),
        ("MLP achieves best AUC (0.875)",
         "Non-linear feature interactions add discriminative value over linear models"),
        ("Logistic Regression best F1 (0.683)",
         "Interpretable linear model remains competitive and is preferred in clinical contexts"),
        ("Glucose + BMI are top predictors",
         "Consistent across ANOVA, RF importance, and clinical literature"),
        ("SMOTE + selection + tuning matter",
         "Each preprocessing step contributed measurable improvement to final performance"),
    ]
    for i, (title, desc) in enumerate(findings):
        ty = 1.72 + i * 0.78
        R(sl, 0.25, ty, 6.1, 0.68, RGBColor(0x25, 0x45, 0x78))
        R(sl, 0.25, ty, 0.055, 0.68, ORANGE)
        T(sl, title, 0.38, ty+0.04, 5.84, 0.30, size=11, bold=True, color=ORANGE)
        T(sl, desc,  0.38, ty+0.33, 5.84, 0.30, size=10, color=LIGHT_BLUE)

    # Right: limitations + future work
    T(sl, "Limitations", 6.65, 1.22, 6.4, 0.42, size=14, bold=True, color=RED)
    lims = [
        "Small dataset (n=768) — Pima women only, poor external validity",
        "High missingness in Insulin (49%) limits feature reliability",
        "No prospective or multi-centre validation performed",
        "Temporal bias: dataset from 1988, population may have shifted",
    ]
    for i, l in enumerate(lims):
        R(sl, 6.65, 1.72+i*0.58, 6.43, 0.50, RGBColor(0x3A, 0x1A, 0x1A))
        T(sl, f"⚠  {l}", 6.80, 1.75+i*0.58, 6.13, 0.42, size=11, color=RGBColor(0xFF, 0xC0, 0xB0))

    T(sl, "Future Work", 6.65, 4.12, 6.4, 0.42, size=14, bold=True, color=GREEN)
    future_items = [
        "Validate on larger, diverse multi-ethnic cohorts",
        "Add HbA1c and fasting glucose as additional features",
        "Implement calibration curves for probability reliability",
        "Fairness audit across age, sex, and ethnic subgroups",
        "Deploy as REST API with explainability (SHAP) layer",
    ]
    for i, f in enumerate(future_items):
        R(sl, 6.65, 4.65+i*0.52, 6.43, 0.46, RGBColor(0x1A, 0x2F, 0x1A))
        T(sl, f"→  {f}", 6.80, 4.67+i*0.52, 6.13, 0.40,
          size=11, color=RGBColor(0xB0, 0xFF, 0xB0))

    # Footer CTA
    T(sl, "Thank you  |  Questions?",
      0.25, 6.60, 7.0, 0.42, size=16, bold=True, color=WHITE)
    T(sl, "Code: github.com/JamshaidAmjad/AI-Diabetes-Prediction",
      0.25, 7.04, 8.0, 0.32, size=11, color=LIGHT_BLUE)

# ═══════════════════════════════════════════════════════════════════════════════
# BUILD
# ═══════════════════════════════════════════════════════════════════════════════
def build():
    prs = _prs()
    s01_title(prs)
    s02_agenda(prs)
    s03_problem(prs)
    s04_dataset(prs)
    s05_eda1(prs)
    s06_eda2(prs)
    s07_pipeline(prs)
    s08_features(prs)
    s09_lr(prs)
    s10_rf(prs)
    s11_mlp(prs)
    s12_training(prs)
    s13_results(prs)
    s14_comparison(prs)
    s15_conclusion(prs)

    out = os.path.join(OUT, "presentation.pptx")
    prs.save(out)
    print(f"Saved {len(prs.slides)}-slide presentation → {out}")
    return out

if __name__ == "__main__":
    build()
